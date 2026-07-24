#!/usr/bin/env python3
"""Add the smaller vacuum-chamber support stand to the system-overview schematic.

Requested by R. Guymon (PR #12, 2026-07-23/24): "include in fig 1 a smaller
support stand under the vacuum chamber, and reflect this as needed in the
text." This reverses the earlier PR #3 correction that removed the
chamber-side stand ("there is no support stand for the vacuum chamber") --
per R. Guymon there IS a smaller stand under the vacuum chamber; the ceiling
cables still suspend the pyrometer housing.

The stand is drawn like the existing heating-head support (gray post on a
black foot bar), but shorter and thinner, rising from the floor line to the
underside of the KF40 cross at the base of the vacuum column. It is inserted
early in the shape tree so the bellows hoses and relief valve render on top
of it, matching how the heating-head post sits behind the bellows.

Edits docs/induction-furnace-schematic-v2.pptx in place (the file produced by
paper/make_schematic_v2.py on PR #3's branch), then re-renders
paper/figures/fig_system_overview_v2.png via LibreOffice + pdftoppm at
300 dpi with the same autocrop as paper/build_schematic_figures.py.

Usage::

    sudo apt-get install -y libreoffice-impress poppler-utils
    pip install python-pptx pillow
    python3 paper/add_chamber_stand_schematic.py
"""
import os
import shutil
import subprocess
import sys
import tempfile

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

REPO_ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
PPTX = os.path.join(REPO_ROOT, "docs", "induction-furnace-schematic-v2.pptx")
OUT_PNG = os.path.join(REPO_ROOT, "paper", "figures", "fig_system_overview_v2.png")
DPI = "300"

GRAY = RGBColor(0xA6, 0xA6, 0xA6)
BLACK = RGBColor(0, 0, 0)


def add_stand() -> None:
    prs = Presentation(PPTX)
    slide = prs.slides[0]
    shapes = slide.shapes

    # Idempotence: the marker text box below is only added by this script.
    for sh in shapes:
        if sh.has_text_frame and sh.text_frame.text == "chamber-stand-marker":
            print("chamber stand already present; nothing to do")
            return

    # Post under the KF40 cross (cross bottom is at y = 8.30 in, spanning
    # x = 8.39-8.76 in); floor line matches the head-support foot (y = 9.435).
    # x = 8.44-8.55 keeps clear of the relief-valve stub at x = 8.71.
    post = shapes.add_shape(1, Inches(8.44), Inches(8.30), Inches(0.11), Inches(1.135))
    post.fill.solid()
    post.fill.fore_color.rgb = GRAY
    post.line.color.rgb = BLACK
    post.line.width = Pt(1.0)
    post.shadow.inherit = False
    sp_tree = post._element.getparent()
    sp_tree.remove(post._element)
    sp_tree.insert(2, post._element)  # render behind the bellows/valve

    foot = shapes.add_shape(1, Inches(8.26), Inches(9.435), Inches(0.50), Inches(0.055))
    foot.fill.solid()
    foot.fill.fore_color.rgb = BLACK
    foot.line.fill.background()
    foot.shadow.inherit = False

    # Invisible marker for idempotence (empty-looking, zero-ish footprint).
    marker = shapes.add_textbox(Inches(0.05), Inches(0.05), Inches(0.3), Inches(0.2))
    marker.text_frame.text = "chamber-stand-marker"
    for para in marker.text_frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(1)
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    prs.save(PPTX)
    print("added chamber support stand to", PPTX)


def render() -> None:
    if not (shutil.which("soffice") and shutil.which("pdftoppm")):
        print("WARNING: soffice/pdftoppm not found; skipping render", file=sys.stderr)
        return
    with tempfile.TemporaryDirectory() as tmp:
        subprocess.run(
            ["soffice", "--headless", "--convert-to", "pdf", "--outdir", tmp, PPTX],
            check=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL,
        )
        pdf = os.path.join(tmp, os.path.splitext(os.path.basename(PPTX))[0] + ".pdf")
        prefix = os.path.join(tmp, "page")
        subprocess.run(
            ["pdftoppm", "-png", "-r", DPI, "-f", "1", "-l", "1", pdf, prefix],
            check=True,
        )
        rendered = next(
            (os.path.join(tmp, f) for f in sorted(os.listdir(tmp))
             if f.startswith("page") and f.endswith(".png")),
            None,
        )
        if rendered is None:
            raise RuntimeError("render produced no PNG")
        shutil.copyfile(rendered, OUT_PNG)
    _autocrop(OUT_PNG)
    print("wrote", OUT_PNG)


def _autocrop(path: str, margin: int = 24) -> None:
    """Same autocrop as paper/build_schematic_figures.py."""
    from PIL import Image, ImageChops
    im = Image.open(path).convert("RGB")
    bg = Image.new("RGB", im.size, (255, 255, 255))
    bbox = ImageChops.difference(im, bg).getbbox()
    if bbox is None:
        return
    left, top, right, bottom = bbox
    im.crop((max(0, left - margin), max(0, top - margin),
             min(im.width, right + margin), min(im.height, bottom + margin))).save(path)


if __name__ == "__main__":
    add_stand()
    render()
