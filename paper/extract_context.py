#!/usr/bin/env python3
"""Extract machine-readable context from the binary source documents in ``docs/``.

The induction-furnace project accumulated context over many years inside binary
office files (``.pptx``, ``.docx``, ``.xlsx``). This utility pulls their text — and
the raster images embedded in the key PowerPoint files — into version-controlled
markdown/figures under ``paper/extracted-context/`` so the HardwareX manuscript can
be drafted and reviewed from text. The binary files in ``docs/`` remain the source
of record.

Usage::

    pip install python-pptx python-docx openpyxl Pillow
    python paper/extract_context.py

Re-running is idempotent: it overwrites the generated files under
``paper/extracted-context/``.
"""
from __future__ import annotations

import os
import sys
import traceback
import zipfile

REPO_ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
DOCS = os.path.join(REPO_ROOT, "docs")
OUT = os.path.join(REPO_ROOT, "paper", "extracted-context")
FIGURES = os.path.join(OUT, "figures")

# Minimum embedded-image size (bytes) worth keeping as a candidate figure;
# smaller assets are spacers/icons/bullets.
MIN_IMAGE_BYTES = 5000

# Key context-rich documents (relative to docs/) -> (output name, description).
TEXT_FILES = [
    ("induction_SOP_200901.docx", "SOP_induction_200901.md",
     "LEPEL Induction Furnace Standard Operating Procedure (Sept 2020)"),
    ("SOP.docx", "SOP_alternate.md", "Alternate Standard Operating Procedure"),
    ("induction_parts_list.xlsx", "parts_list.md",
     "Induction furnace parts list / bill of materials (working)"),
    ("induction-furnace-schematic.pptx", "schematic_induction_furnace.md",
     "Induction furnace system schematic (component callouts)"),
    ("schematic.pptx", "schematic_support_stand.md",
     "System and support-stand schematic"),
    ("induction_order_corrections.pptx", "order_corrections_coils.md",
     "Coil order corrections (geometry / specifications)"),
    ("manual/useful_links.docx", "useful_links.md",
     "Useful links (oscillator/rectifier tubes)"),
]

# PowerPoint files whose embedded raster images are worth extracting as figures.
IMAGE_SOURCES = [
    ("induction-furnace-schematic.pptx", "induction-furnace-schematic"),
    ("schematic.pptx", "system-schematic"),
    ("induction_order_corrections.pptx", "coil-order-corrections"),
]


def extract_pptx(path: str) -> str:
    from pptx import Presentation

    out = []
    for i, slide in enumerate(Presentation(path).slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                texts.append(shape.text_frame.text.strip())
            if shape.has_table:
                for row in shape.table.rows:
                    texts.append(" | ".join(c.text.strip() for c in row.cells))
        body = "\n".join(texts) if texts else "_(no text; image/diagram slide)_"
        out.append(f"### Slide {i}\n\n{body}")
    return "\n\n".join(out)


def extract_docx(path: str) -> str:
    from docx import Document

    doc = Document(path)
    out = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    for ti, tbl in enumerate(doc.tables, 1):
        out.append(f"\n**Table {ti}:**\n")
        for row in tbl.rows:
            out.append("| " + " | ".join(c.text.strip() for c in row.cells) + " |")
    return "\n\n".join(out)


def extract_xlsx(path: str) -> str:
    import openpyxl

    wb = openpyxl.load_workbook(path, data_only=True, read_only=True)
    out = []
    for ws in wb.worksheets:
        out.append(f"### Sheet: {ws.title}\n")
        rows = 0
        for row in ws.iter_rows(values_only=True):
            vals = [str(c) if c is not None else "" for c in row]
            if any(v.strip() for v in vals):
                out.append("| " + " | ".join(vals) + " |")
                rows += 1
            if rows > 300:
                out.append("_... (truncated)_")
                break
    return "\n".join(out)


EXTRACTORS = {"pptx": extract_pptx, "docx": extract_docx, "xlsx": extract_xlsx}


def extract_text_files() -> list[str]:
    os.makedirs(OUT, exist_ok=True)
    index = [
        "# Extracted context from binary documents",
        "",
        "Machine-readable text extracted from the binary source documents in "
        "[`../../docs/`](../../docs/) so the HardwareX manuscript can be drafted and "
        "reviewed from version-controlled text. Each file links back to its binary "
        "source of record. Regenerate with `python paper/extract_context.py`.",
        "",
        "| Extracted file | Source document | Description |",
        "|----------------|-----------------|-------------|",
    ]
    for rel, outname, desc in TEXT_FILES:
        src = os.path.join(DOCS, rel)
        if not os.path.exists(src):
            print(f"SKIP (missing): {rel}", file=sys.stderr)
            continue
        ext = rel.lower().rsplit(".", 1)[-1]
        try:
            content = EXTRACTORS[ext](src)
        except Exception:  # pragma: no cover - defensive logging
            print(f"ERROR extracting {rel}", file=sys.stderr)
            traceback.print_exc()
            continue
        header = (
            f"# {desc}\n\n"
            f"> Extracted from [`docs/{rel}`](../../docs/{rel}). Auto-extracted text "
            f"for reference; the binary file is the source of record.\n"
        )
        with open(os.path.join(OUT, outname), "w") as f:
            f.write(header + "\n" + content + "\n")
        index.append(f"| [`{outname}`]({outname}) | `docs/{rel}` | {desc} |")
        print(f"WROTE: paper/extracted-context/{outname}")
    with open(os.path.join(OUT, "README.md"), "w") as f:
        f.write("\n".join(index) + "\n")
    print("WROTE: paper/extracted-context/README.md")
    return index


def extract_images() -> None:
    os.makedirs(FIGURES, exist_ok=True)
    manifest = [
        "# Extracted figures",
        "",
        "Raster images embedded in the PowerPoint source documents, pulled out as "
        "candidate manuscript figures (hardware photos, schematics, coil geometry). "
        "Tiny spacer/icon images and vector overlays are omitted. Regenerate with "
        "`python paper/extract_context.py`.",
        "",
        "| Figure | Source document | Original embedded name | Size (bytes) |",
        "|--------|-----------------|------------------------|-------------:|",
    ]
    for rel, prefix in IMAGE_SOURCES:
        src = os.path.join(DOCS, rel)
        if not os.path.exists(src):
            print(f"SKIP (missing): {rel}", file=sys.stderr)
            continue
        with zipfile.ZipFile(src) as z:
            media = sorted(
                n for n in z.namelist()
                if n.startswith("ppt/media/")
                and n.lower().rsplit(".", 1)[-1] in ("png", "jpeg", "jpg")
            )
            idx = 1
            for name in media:
                data = z.read(name)
                if len(data) < MIN_IMAGE_BYTES:
                    continue
                ext = name.lower().rsplit(".", 1)[-1]
                ext = "jpg" if ext == "jpeg" else ext
                outname = f"{prefix}_{idx:02d}.{ext}"
                with open(os.path.join(FIGURES, outname), "wb") as f:
                    f.write(data)
                manifest.append(
                    f"| `{outname}` | `docs/{rel}` | `{os.path.basename(name)}` | {len(data)} |"
                )
                print(f"WROTE: paper/extracted-context/figures/{outname} ({len(data)} bytes)")
                idx += 1
    with open(os.path.join(FIGURES, "README.md"), "w") as f:
        f.write("\n".join(manifest) + "\n")
    print("WROTE: paper/extracted-context/figures/README.md")


def main() -> None:
    extract_text_files()
    extract_images()


if __name__ == "__main__":
    main()
