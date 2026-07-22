#!/usr/bin/env python3
"""Build docs/induction-furnace-schematic-v2.pptx from the v1 schematic.

Requested revisions (R. Guymon, PR #3, 2026-07-07):
  - drop the control-computer icon/label
  - show chiller water entering the heating head (coil is water-cooled)
  - make the right vertical support line fully solid (no dashed top section)
  - keep every word on one line (no mid-word wraps)
  - mount the heating head on a support
  - label the vacuum chamber stack
  - shrink the chiller and generator

Follow-up correction (R. Guymon, PR #3, 2026-07-07): there is no support stand
for the vacuum chamber -- the chamber stack is joined by KF40 flanges to the
pyrometer housing, which is suspended from the ceiling by cables. The stand
supports only the heating head. So: the chamber-side stand/clamps are removed,
both ceiling cables are drawn solid (load-bearing) down to the pyrometer
housing, and the stand label moves under the heating-head support.

Second follow-up (R. Guymon, PR #3, 2026-07-07):
  - cooling water runs chiller -> heating head -> generator -> chiller
  - the 0.5 psi overpressure relief valve hangs from the bottom of the chamber
    stack, at the KF40 cross where the bellows connects

Third follow-up (R. Guymon, PR #3, 2026-07-07):
  - draw only TWO ceiling cables (the housing physically hangs from three, but
    two reads better in this flat side view -- see AGENTS.md)
  - move the generator and chiller in next to the heating head / pump instead
    of leaving them far off at the left edge

Fourth follow-up (R. Guymon, PR #3, 2026-07-08):
  - greatly shorten the ceiling cables (the long run wasted vertical space);
    the ceiling bar now sits just above the pyrometer
  - cooling water verified to run chiller -> generator -> heating head ->
    chiller (reverses the previous head-first order)
  - no specific equipment model names in the drawing ("Turbo Pumping Station"
    instead of "T-Station 85")

Fifth follow-up (R. Guymon, PR #3, 2026-07-13):
  - the coolant is ethylene glycol, not water: the loop label reads
    "Ethylene Glycol Coolant" and the chiller box "Recirculating Chiller"
"""
import copy, os, shutil
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

REPO_ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
SRC = os.path.join(REPO_ROOT, "docs", "induction-furnace-schematic.pptx")
DST = os.path.join(REPO_ROOT, "docs", "induction-furnace-schematic-v2.pptx")
shutil.copyfile(SRC, DST)

prs = Presentation(DST)
slide = prs.slides[0]
shapes = slide.shapes
by_id = {sh.shape_id: sh for sh in shapes}

def delete(sid):
    el = by_id[sid]._element
    el.getparent().remove(el)

def set_text(sh, lines, size, bold=False):
    tf = sh.text_frame
    tf.word_wrap = False
    tf.text = lines[0]
    for extra in lines[1:]:
        tf.add_paragraph().text = extra
    for para in tf.paragraphs:
        para.alignment = PP_ALIGN.CENTER
        for run in para.runs:
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = RGBColor(0, 0, 0)

def nowrap(sid, size=None):
    sh = by_id[sid]
    sh.text_frame.word_wrap = False
    if size is not None:
        for para in sh.text_frame.paragraphs:
            for run in para.runs:
                run.font.size = Pt(size)

def tail_arrow(sh, w="med", ln_="med"):
    ln = sh.line._get_or_add_ln()
    te = ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': w, 'len': ln_})
    ln.append(te)

def style_line(sh, width_pt, rgb=RGBColor(0,0,0), dash=None):
    sh.line.color.rgb = rgb
    sh.line.width = Pt(width_pt)
    if dash is not None:
        ln = sh.line._get_or_add_ln()
        d = ln.makeelement(qn('a:prstDash'), {'val': dash})
        ln.append(d)

def label(text, l, t, w, size=10, align=PP_ALIGN.LEFT, rgb=None):
    tb = shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(0.26))
    tf = tb.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.text = text
    p = tf.paragraphs[0]
    p.alignment = align
    for run in p.runs:
        run.font.size = Pt(size)
        if rgb is not None:
            run.font.color.rgb = rgb
    return tb

def leader(x1, y1, x2, y2):
    c = shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    style_line(c, 1.0)
    tail_arrow(c, w="sm", ln_="sm")
    return c

# ---- 1. remove control computer icon + label, old chiller<->gen arrows, old RF leads
for sid in (9, 137, 121, 122, 135, 136):
    delete(sid)

# ---- 2. all ceiling cables solid: they are load-bearing (the chamber stack
#         hangs from the pyrometer housing, which the cables tie to the ceiling)
for sid in (34, 35):
    ln = by_id[sid]._element.find('.//' + qn('a:ln'))
    for d in ln.findall(qn('a:prstDash')):
        ln.remove(d)

# (Only two cables are drawn. The housing physically hangs from three, but in
#  this flat side view the third attachment sits behind the pyrometer and an
#  angled third line read poorly, so it is deliberately omitted -- R. Guymon,
#  PR #3, 2026-07-07.)

# ---- 2b. greatly shorten the cables (R. Guymon, PR #3, 2026-07-08): drop the
#          ceiling bar to just above the pyrometer so the cables read as short
#          stubs instead of a full-height run that wastes figure space.
CEILING_Y = 4.05
ceiling = by_id[50]
ceiling.top = Inches(CEILING_Y)
for sid in (34, 35):
    cab = by_id[sid]
    bottom = cab.top + cab.height
    cab.top = Inches(CEILING_Y)
    cab.height = Emu(bottom - Inches(CEILING_Y))
lab38 = by_id[38]                 # move the cables label down next to them
lab38.top, lab38.height = Inches(4.30), Inches(0.70)

# ---- 3. shrink + reposition chiller and generator, tucked in next to the
#         heating head / pumping station instead of far off at the left edge
ch = by_id[71]
ch.left, ch.top, ch.width, ch.height = Inches(3.90), Inches(8.50), Inches(1.50), Inches(0.95)
set_text(ch, ["Recirculating", "Chiller"], 11)
gen = by_id[72]
gen.left, gen.top, gen.width, gen.height = Inches(4.00), Inches(6.65), Inches(1.40), Inches(0.90)
set_text(gen, ["Induction", "Generator"], 11)

# ---- 4. RF leads: generator -> heating head (dotted, +/- polarity labels kept)
for y in (6.90, 7.18):
    c = shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(5.40), Inches(y), Inches(6.82), Inches(y))
    style_line(c, 2.0, dash='sysDot')
plus, minus = by_id[133], by_id[134]
plus.left, plus.top = Inches(5.44), Inches(6.53)
minus.left, minus.top = Inches(5.44), Inches(6.91)

# ---- 5. coolant loop (ethylene glycol, not water -- R. Guymon, PR #3,
#         2026-07-13): chiller -> generator -> heating head -> chiller
#         (order verified by R. Guymon, PR #3, 2026-07-08)
BLUE = RGBColor(0x1F, 0x6F, 0xC4)
def coolant(points):
    fb = shapes.build_freeform(Emu(Inches(points[0][0])), Emu(Inches(points[0][1])), scale=1.0)
    fb.add_line_segments([(Emu(Inches(x)), Emu(Inches(y))) for x, y in points[1:]], close=False)
    sh = fb.convert_to_shape()
    sh.fill.background()
    sh.shadow.inherit = False
    style_line(sh, 2.0, rgb=BLUE)
    tail_arrow(sh)
    return sh
coolant([(4.30, 8.50), (4.30, 7.55)])                                             # chiller -> generator
coolant([(4.90, 7.55), (4.90, 7.65), (6.95, 7.65), (6.95, 7.31)])                 # generator -> head
coolant([(7.20, 7.31), (7.20, 7.75), (5.57, 7.75), (5.57, 8.70), (5.42, 8.70)])   # head -> chiller
# label sits left of the chiller->generator run (right-aligned so its right
# edge stops just short of the vertical coolant line at x=4.30)
label("Ethylene Glycol Coolant", 2.40, 7.98, 1.85, size=10,
      align=PP_ALIGN.RIGHT, rgb=BLUE)

# ---- 6. heating-head support post standing on its own base bar; the chamber
#         stack itself has NO stand (it hangs from the ceiling cables), so the
#         chamber-side stand/clamps group is removed entirely
post = shapes.add_shape(1, Inches(7.38), Inches(7.29), Inches(0.13), Inches(2.15))  # 1 = rectangle
post.fill.solid(); post.fill.fore_color.rgb = RGBColor(0xA6, 0xA6, 0xA6)
post.line.color.rgb = RGBColor(0, 0, 0); post.line.width = Pt(1.0)
post.shadow.inherit = False
spTree = post._element.getparent()
spTree.remove(post._element)
spTree.insert(2, post._element)   # first drawable -> everything else renders on top
foot = shapes.add_shape(1, Inches(7.28), Inches(9.435), Inches(0.90), Inches(0.055))
foot.fill.solid(); foot.fill.fore_color.rgb = RGBColor(0, 0, 0)
foot.line.fill.background()
foot.shadow.inherit = False
delete(115)                       # chamber support stand + clamps: does not exist
base = by_id[91]                  # relabel + center the label under the head support
base.left, base.width = Inches(6.93), Inches(1.60)
set_text(base, ["Support Stand"], 12)

# ---- 7. bellows text off the hose (support post would cross it)
by_id[14].text_frame.text = ""
label("Bellows", 5.72, 8.28, 0.70, size=10)

# ---- 8. single-line words everywhere; no equipment model names in the
#         drawing (R. Guymon, PR #3, 2026-07-08)
set_text(by_id[11], ["Turbo Pumping", "Station"], 11)
nowrap(23, 12)          # Pyrometer (rotated)
nowrap(24, 9)           # Pressure Sensor
set_text(by_id[25], ["Pyrometer Housing"], 12)   # matches manuscript terminology
nowrap(26, 12)          # Quartz Disc
set_text(by_id[38], ["Fixturing Cables", "(to ceiling)"], 12)
set_text(by_id[70], ["Heating", "Head"], 11)

# ---- 9. crucible inside the quartz tube, at coil height
cru = shapes.add_shape(1, Inches(8.47), Inches(6.95), Inches(0.20), Inches(0.15))
cru.fill.solid(); cru.fill.fore_color.rgb = RGBColor(0x59, 0x59, 0x59)
cru.line.color.rgb = RGBColor(0, 0, 0); cru.line.width = Pt(0.75)
cru.shadow.inherit = False

# ---- 10. vacuum-stack labels with leader arrows
# (label sits left of the repositioned generator; its leader passes over the
#  generator/head gap to reach the tube)
label("Quartz Tube (Vacuum Chamber)", 3.17, 6.33, 2.18, align=PP_ALIGN.RIGHT)
leader(5.42, 6.46, 8.40, 6.55)
label("KF40 Fitting", 9.02, 6.16, 0.95)
leader(9.00, 6.29, 8.74, 6.29)
label("Crucible", 9.05, 6.52, 0.85)
leader(9.03, 6.65, 8.66, 6.98)
label("Induction Coil", 9.02, 6.93, 0.95)
leader(9.00, 7.06, 8.77, 7.06)
label("KF40 Cross", 9.02, 7.80, 0.95)
leader(9.00, 7.93, 8.74, 7.93)

# ---- 11. overpressure relief valve (0.5 psi cracking pressure) hanging from
#          the bottom of the chamber stack, at the KF40 cross where the bellows
#          connects (matches Fig. fig_vacuum_details(b) and the vacuum-path text)
stub = shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(8.71), Inches(8.30), Inches(8.71), Inches(8.55))
style_line(stub, 2.0)
rv = shapes.add_shape(1, Inches(8.63), Inches(8.55), Inches(0.16), Inches(0.14))
rv.fill.solid(); rv.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
rv.line.color.rgb = RGBColor(0, 0, 0); rv.line.width = Pt(1.25)
rv.shadow.inherit = False
label("Relief Valve", 9.02, 8.49, 0.95)
leader(9.00, 8.62, 8.81, 8.62)

prs.save(DST)
print("saved", DST)
