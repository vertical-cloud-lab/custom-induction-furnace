#!/usr/bin/env python3
"""Build docs/ysz-stack-schematic.pptx: the YSZ sample assembly, drawn in the
same style as the v2 system-overview schematic (make_schematic_v2.py).

Requested by R. Guymon (PR #3, 2026-07-07) from a hand sketch: the fresh quartz
tube (~35 mm ID) with the induction coil turns outside it, and inside, top to
bottom, a tantalum susceptor block / YSZ specimen / tantalum susceptor block
(25.5 mm) sandwich resting on a ceramic stub (28 mm; boron nitride for best
results, sometimes MgO --- a solid stub, not a crucible, per R. Guymon PR #3
2026-07-10), which sits on the alumina support tube that passes down to the
KF40 hardware at the bottom of the vacuum column.

Corrections (S. Baird, PR #12, 2026-07-22): the BN stub is a diffusion
barrier (it prevents reaction between the tantalum and the alumina), not heat
dissipation; the alumina support is a tube, not a rod, with holes bored
laterally along its outer perimeter to aid evacuation; and a teflon tube sits
between the alumina tube and the KF40 fitting.

Correction (R. Guymon, PR #12, 2026-07-23): the teflon tube is underneath the
alumina tube -- the alumina tube rests on top of the teflon tube, it does not
sit inside it.

If LibreOffice (soffice) and pdftoppm are available, this script also renders
the slide to paper/figures/fig_ysz_stack.png (300 dpi, autocropped) the same
way build_schematic_figures.py does on PR #3's branch.
"""
import os
import shutil
import subprocess
import sys
import tempfile
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

REPO_ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
DST = os.path.join(REPO_ROOT, "docs", "ysz-stack-schematic.pptx")

prs = Presentation()
prs.slide_width = Inches(7.0)
prs.slide_height = Inches(9.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
shapes = slide.shapes

BLACK = RGBColor(0, 0, 0)
GRAY = RGBColor(0xD9, 0xD9, 0xD9)     # tantalum blocks
RED = RGBColor(0xC0, 0x00, 0x00)      # YSZ specimen (matches the sketch)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# Geometry is to scale horizontally: 1 mm = 0.0314 in (tube ID 35 mm = 1.10 in).
CX = 3.30                 # tube centerline
WALL_L, WALL_R = 2.75, 3.85


def style_line(sh, width_pt, rgb=BLACK, dash=None):
    sh.line.color.rgb = rgb
    sh.line.width = Pt(width_pt)
    if dash is not None:
        ln = sh.line._get_or_add_ln()
        d = ln.makeelement(qn('a:prstDash'), {'val': dash})
        ln.append(d)


def line(x1, y1, x2, y2, width_pt=2.0):
    c = shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                             Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    style_line(c, width_pt)
    c.shadow.inherit = False
    return c


def box(l, t, w, h, fill=WHITE, line_pt=1.25):
    sh = shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))  # 1 = rectangle
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = BLACK
    sh.line.width = Pt(line_pt)
    sh.shadow.inherit = False
    return sh


def oval(l, t, w, h, line_pt=1.5):
    sh = shapes.add_shape(9, Inches(l), Inches(t), Inches(w), Inches(h))  # 9 = oval
    sh.fill.solid()
    sh.fill.fore_color.rgb = WHITE
    sh.line.color.rgb = BLACK
    sh.line.width = Pt(line_pt)
    sh.shadow.inherit = False
    return sh


def label(text, l, t, w, size=10, align=PP_ALIGN.LEFT, rgb=None):
    tb = shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(0.26))
    tf = tb.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.text = text
    p = tf.paragraphs[0]
    p.alignment = align
    for run in p.runs:
        run.font.size = Pt(size)
        if rgb is not None:
            run.font.color.rgb = rgb
    return tb


def leader(x1, y1, x2, y2):
    c = shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                             Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    style_line(c, 1.0)
    c.shadow.inherit = False
    ln = c.line._get_or_add_ln()
    te = ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'sm', 'len': 'sm'})
    ln.append(te)
    return c


def dimension(x1, x2, y, text, text_above=True, size=9):
    """Horizontal dimension line with end ticks and a centered value label."""
    line(x1, y, x2, y, 1.0)
    for x in (x1, x2):
        line(x, y - 0.05, x, y + 0.05, 1.0)
    ty = y - 0.24 if text_above else y + 0.06
    label(text, x1, ty, x2 - x1, size=size, align=PP_ALIGN.CENTER)


# ---- KF40 hardware at the bottom of the vacuum column (drawn first so the
#      teflon tube renders on top of it, visibly seated inside)
box(2.60, 7.35, 1.40, 0.35, line_pt=1.5)          # KF40 fitting at the tube base
box(2.55, 7.75, 1.50, 0.55, line_pt=1.5)          # KF40 cross below it

# ---- teflon tube seated in the KF40 fitting; the alumina support tube rests
#      on its top rim (R. Guymon, PR #12, 2026-07-23)
box(2.95, 7.15, 0.70, 1.00, line_pt=1.25)

# ---- alumina support tube (below the susceptor blocks, resting ON TOP of
#      the teflon tube -- its bottom edge coincides with the teflon top rim)
box(3.05, 4.17, 0.50, 2.98)
# hollow bore, shown as hidden (dashed) lines
for bx in (3.19, 3.41):
    c = shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                             Inches(bx), Inches(4.22), Inches(bx), Inches(7.10))
    style_line(c, 0.75, dash='dash')
    c.shadow.inherit = False
# lateral evacuation holes bored along the outer perimeter (side view: one
# small hole on each silhouette edge at several heights)
for hy in (4.70, 5.45, 6.20, 6.95):
    for hx in (3.05, 3.55):
        oval(hx - 0.045, hy - 0.045, 0.09, 0.09, line_pt=1.0)

# ---- quartz tube walls (vacuum chamber), ~35 mm ID
line(WALL_L, 0.95, WALL_L, 7.35)
line(WALL_R, 0.95, WALL_R, 7.35)
dimension(WALL_L, WALL_R, 0.72, "~35 mm (tube ID)")

# ---- induction coil turns outside the tube (cross-section: one circle per turn)
for cy in (1.85, 2.25, 2.65, 3.05):
    oval(2.39, cy - 0.16, 0.32, 0.32)
    oval(3.89, cy - 0.16, 0.32, 0.32)

# ---- BN diffusion-barrier stub (28 mm): a solid ceramic stub, not a
#      crucible; it separates the tantalum from the alumina tube below.
#      The Ta/YSZ/Ta sandwich rests on its flat top (dimensioned below the
#      stub so the line does not read as a division of the solid block)
box(2.86, 3.21, 0.88, 0.96, line_pt=1.5)
dimension(2.86, 3.74, 4.40, "28 mm", text_above=False)

# ---- tantalum / YSZ / tantalum sandwich (25.5 mm blocks); the lower block
#      rests on top of the stub
box(2.90, 1.82, 0.80, 0.60, fill=GRAY)            # upper Ta susceptor block
ysz = box(2.92, 2.42, 0.76, 0.07, fill=RED, line_pt=0.75)  # YSZ specimen
ysz.line.color.rgb = RED
box(2.90, 2.49, 0.80, 0.72, fill=GRAY)            # lower Ta susceptor block
dimension(2.90, 3.70, 1.62, "25.5 mm")
label("Tantalum", 2.90, 2.03, 0.80, size=9, align=PP_ALIGN.CENTER)
label("Tantalum", 2.90, 2.76, 0.80, size=9, align=PP_ALIGN.CENTER)

# ---- labels with leader arrows (v2 style: left/right of the stack)
label("Quartz Tube (Vacuum Chamber)", 0.28, 1.02, 2.10, align=PP_ALIGN.RIGHT)
leader(2.40, 1.15, 2.72, 1.15)
label("Induction Coil", 1.05, 2.53, 1.05, align=PP_ALIGN.RIGHT)
leader(2.14, 2.66, 2.37, 2.66)

# the YSZ leader threads the gap between the second and third coil turns
label("YSZ Specimen", 4.72, 2.35, 1.20)
leader(4.68, 2.455, 3.70, 2.455)
label("BN Diffusion Barrier", 4.72, 3.48, 2.00)
leader(4.68, 3.60, 3.78, 3.60)
label("Alumina Support Tube", 4.72, 5.22, 1.70)
label("(lateral evacuation holes)", 4.72, 5.44, 1.90, size=9)
leader(4.68, 5.34, 3.59, 5.34)
label("Teflon Tube", 4.72, 6.95, 1.00)
leader(4.68, 7.07, 3.62, 7.28)
label("KF40 Fitting", 4.72, 7.40, 1.00)
leader(4.68, 7.52, 4.04, 7.52)
label("KF40 Cross", 4.72, 7.92, 1.00)
leader(4.68, 8.03, 4.09, 8.03)

prs.save(DST)
print("saved", DST)


# ---- optional render to paper/figures/fig_ysz_stack.png (matches the
#      build_schematic_figures.py pipeline: soffice -> pdf -> pdftoppm 300 dpi
#      -> autocrop the white slide margins)
def _autocrop(path, margin=24):
    from PIL import Image, ImageChops
    im = Image.open(path).convert("RGB")
    bg = Image.new("RGB", im.size, (255, 255, 255))
    bbox = ImageChops.difference(im, bg).getbbox()
    if bbox is None:
        return
    left, top, right, bottom = bbox
    im.crop((max(0, left - margin), max(0, top - margin),
             min(im.width, right + margin),
             min(im.height, bottom + margin))).save(path)


if shutil.which("soffice") and shutil.which("pdftoppm"):
    out_png = os.path.join(REPO_ROOT, "paper", "figures", "fig_ysz_stack.png")
    with tempfile.TemporaryDirectory() as tmp:
        subprocess.run(["soffice", "--headless", "--convert-to", "pdf",
                        "--outdir", tmp, DST], check=True,
                       stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
        pdf = os.path.join(tmp, "ysz-stack-schematic.pdf")
        subprocess.run(["pdftoppm", "-png", "-r", "300", "-f", "1", "-l", "1",
                        pdf, os.path.join(tmp, "page")], check=True)
        rendered = next(os.path.join(tmp, f) for f in sorted(os.listdir(tmp))
                        if f.startswith("page") and f.endswith(".png"))
        shutil.copyfile(rendered, out_png)
    _autocrop(out_png)
    print("rendered", out_png)
else:
    print("WARNING: soffice/pdftoppm not found; fig_ysz_stack.png not "
          "re-rendered", file=sys.stderr)
